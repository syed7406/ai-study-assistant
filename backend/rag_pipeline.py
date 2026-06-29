from __future__ import annotations

import os
from pathlib import Path
import uuid
from dataclasses import dataclass
from functools import lru_cache
from typing import Any, Dict, List, Tuple

from docx import Document
import fitz  # PyMuPDF
import numpy as np
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from vector_store import add_documents, query_documents


MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"


@dataclass
class ChunkResult:
    chunk_id: str
    text: str
    metadata: Dict[str, Any]


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}


def _safe_int(value: Any, default: int) -> int:
    try:
        return int(value)
    except Exception:
        return default


def _split_text(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
    cleaned = "\n".join(line.rstrip() for line in text.splitlines()).strip()
    if not cleaned:
        return []

    chunks: List[str] = []
    start = 0
    while start < len(cleaned):
        end = min(start + chunk_size, len(cleaned))
        chunk = cleaned[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(cleaned):
            break
        start = max(end - chunk_overlap, start + 1)
    return chunks


def _extract_pdf(file_bytes: bytes) -> List[tuple[int, str]]:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages: List[tuple[int, str]] = []
    for page_index in range(len(doc)):
        text = doc[page_index].get_text("text").strip()
        if text:
            pages.append((page_index + 1, text))
    return pages


def _extract_pptx(file_bytes: bytes) -> List[tuple[int, str]]:
    import io

    presentation = Presentation(io.BytesIO(file_bytes))
    slides: List[tuple[int, str]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        slide_text = "\n".join(parts).strip()
        if slide_text:
            slides.append((slide_index, slide_text))
    return slides


def _extract_docx(file_bytes: bytes) -> List[tuple[int, str]]:
    import io

    document = Document(io.BytesIO(file_bytes))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    text = "\n".join(paragraphs).strip()
    return [(1, text)] if text else []


def _extract_txt_or_md(file_bytes: bytes) -> List[tuple[int, str]]:
    text = ""
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = file_bytes.decode(encoding).strip()
            break
        except UnicodeDecodeError:
            continue
    return [(1, text)] if text else []


def _extract_by_extension(file_bytes: bytes, filename: str) -> List[tuple[int, str]]:
    extension = Path(filename).suffix.lower()
    if extension == ".pdf":
        return _extract_pdf(file_bytes)
    if extension == ".pptx":
        return _extract_pptx(file_bytes)
    if extension == ".docx":
        return _extract_docx(file_bytes)
    if extension in {".txt", ".md"}:
        return _extract_txt_or_md(file_bytes)
    raise ValueError(
        f"Unsupported file type '{extension}'. Supported types: PDF, PPTX, DOCX, TXT, MD."
    )


@dataclass
class RAGConfig:
    chunk_size: int = _safe_int(os.getenv("MAX_CHUNK_SIZE", "500"), 500)
    chunk_overlap: int = _safe_int(os.getenv("CHUNK_OVERLAP", "50"), 50)
    top_k: int = _safe_int(os.getenv("TOP_K_RESULTS", "5"), 5)
    min_similarity: float = float(os.getenv("MIN_SIMILARITY", "0.25"))


@lru_cache(maxsize=1)
def get_embedding_model() -> SentenceTransformer:
    return SentenceTransformer(MODEL_NAME)


def extract_document_chunks(file_bytes: bytes, filename: str, chunk_size: int, chunk_overlap: int) -> List[ChunkResult]:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{extension}'. Supported types: PDF, PPTX, DOCX, TXT, MD."
        )

    sections = _extract_by_extension(file_bytes, filename)
    chunks: List[ChunkResult] = []

    for section_number, text in sections:
        for chunk_index, chunk_text in enumerate(_split_text(text, chunk_size, chunk_overlap), start=1):
            chunks.append(
                ChunkResult(
                    chunk_id=str(uuid.uuid4()),
                    text=chunk_text,
                    metadata={
                        "filename": filename,
                        "section": section_number,
                        "chunk": chunk_index,
                        "file_type": extension.lstrip("."),
                    },
                )
            )

    return chunks


def embed_texts(texts: List[str]) -> List[List[float]]:
    model = get_embedding_model()
    embeddings = model.encode(texts, normalize_embeddings=True)
    return embeddings.tolist()


def ingest_document(user_id: str, file_bytes: bytes, filename: str) -> int:
    config = RAGConfig()
    chunks = extract_document_chunks(file_bytes, filename, config.chunk_size, config.chunk_overlap)
    if not chunks:
        return 0

    embeddings = embed_texts([chunk.text for chunk in chunks])
    add_documents(
        user_id=user_id,
        ids=[chunk.chunk_id for chunk in chunks],
        documents=[chunk.text for chunk in chunks],
        metadatas=[chunk.metadata for chunk in chunks],
        embeddings=embeddings,
    )
    return len(chunks)


def retrieve_relevant_chunks(user_id: str, question: str, top_k: int | None = None):
    config = RAGConfig()
    k = top_k or config.top_k
    question_embedding = embed_texts([question])[0]
    raw = query_documents(user_id, question_embedding, k)

    documents = (raw.get("documents") or [[]])[0]
    metadatas = (raw.get("metadatas") or [[]])[0]
    distances = (raw.get("distances") or [[]])[0]

    filtered_docs: List[str] = []
    filtered_meta: List[Dict[str, Any]] = []
    filtered_dist: List[float] = []

    for doc, meta, distance in zip(documents, metadatas, distances):
        similarity = 1 - float(distance)
        if similarity >= config.min_similarity:
            filtered_docs.append(doc)
            filtered_meta.append(meta)
            filtered_dist.append(float(distance))

    return {
        "documents": [filtered_docs],
        "metadatas": [filtered_meta],
        "distances": [filtered_dist],
    }


def calculate_confidence(distances: List[float]) -> float:
    if not distances:
        return 0.0
    top_distance = float(distances[0])
    similarity = (1 - top_distance) * 100
    return round(float(np.clip(similarity, 0, 100)), 1)
